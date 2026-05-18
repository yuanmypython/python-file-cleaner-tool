import os
from typing import cast
import tkinter
from tkinter import filedialog
import hashlib
from tkinter import messagebox
def get_md5(file_path):
    md5 = hashlib.md5()
    with open(file_path, "rb") as f:
        while True:
            chunk = f.read(8192)
            if not chunk:
                break
            md5.update(chunk)
    return md5.hexdigest()
root = tkinter.Tk()
root.withdraw()
wait = filedialog.askdirectory(title="请选择要处理的文件夹")
if not wait:
    messagebox.showinfo("提示","您已取消选择")
from openpyxl import load_workbook

def is_empty_excel(file_path):
    """判断xlsx文件是否真的空白（所有工作表都没有任何数据）"""
    try:
        # 只读模式打开，速度快，不会占用文件
        wb = load_workbook(file_path, read_only=True, data_only=True)
        # 遍历所有工作表
        for sheet in wb.worksheets:
            # 只要有一个工作表有数据，就不是空白
            if sheet.max_row > 1 or sheet.max_column > 1:
                return False
            # 特殊情况：只有A1单元格有内容
            if sheet.max_row == 1 and sheet.max_column == 1:
                cell_value = sheet.cell(row=1, column=1).value
                if cell_value is not None and str(cast(str,cell_value)).strip() != "": #pyright:ignore
                    return False
        wb.close()
        return True
    except Exception as e:
        # 文件损坏、不是xlsx、被占用等情况，跳过
        print(f"无法检测Excel文件 {file_path}: {e}")
        return False
from docx import Document
from pptx import Presentation

def is_empty_docx(file_path):
    """判断空白.docx文件"""
    try:
        doc = Document(file_path)
        # 检查所有段落和表格
        for para in doc.paragraphs:
            if para.text.strip() != "":
                return False
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip() != "":
                        return False
        return True
    except Exception as e:
        print(f"无法检测docx文件 {file_path}: {e}")
        return False

def is_empty_pptx(file_path):
    """判断空白.pptx文件"""
    try:
        prs = Presentation(file_path)
        # 只要有一页幻灯片有内容，就不是空白
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip() != "":
                    return False
        return True
    except Exception as e:
        print(f"无法检测pptx文件 {file_path}: {e}")
        return False


del_no = 0
del_copy = 0
del_d = 0
md5_dict = {}
to_delete_empty = []
for i_path,i_file,in_file in os.walk(wait,topdown=False):
    for filename in in_file:
        del_filename = os.path.join(i_path,filename)
        file_ext = os.path.splitext(del_filename)[1].lower()
        if os.path.getsize(del_filename) == 0:
            to_delete_empty.append(del_filename)
            continue

        if file_ext == ".xlsx":
            if is_empty_excel(del_filename):
                to_delete_empty.append(del_filename)
                continue
        if file_ext == ".docx":
            if is_empty_docx(del_filename):
                to_delete_empty.append(del_filename)
        if file_ext == ".pptx":
            if is_empty_pptx(del_filename):
                to_delete_empty.append(del_filename)
        if file_ext == '.xls':
            if is_empty_excel(del_filename):
                to_delete_empty.append(del_filename)
if to_delete_empty:
    display_text = ""
    for i,file in enumerate(to_delete_empty[:20]):
        display_text += f"{i+1}. {file}\n"
        if len(to_delete_empty)>20:
            display_text += f"还有{len(to_delete_empty)-20}个，完整列表看控制台"
    confirm = messagebox.askokcancel(
        "确认删除空文件",
        f"共检测到{len(to_delete_empty)}个空文件/空白xlsx:\n\n"
        f"{display_text}\n"
        "是否全部删除？"
    )
    if confirm:
        for file in to_delete_empty:
            try:
                os.remove(file)
                del_no += 1
            except Exception as e:
                print(f"删除失败，{file}原因：{e}")
    else:
        print("已取消删除全部文件")
for i_path, i_file, in_file in os.walk(wait, topdown=False):
    for filename in in_file:
        del_filename = os.path.join(i_path, filename)
        if not os.path.exists(del_filename):
            continue

        md5_value = get_md5(del_filename)
        if md5_value in md5_dict:
            os.remove(del_filename)
            del_copy += 1
            continue
        else:
            md5_dict[md5_value] = del_filename

for i_path,i_file,_in_file in os.walk(wait,topdown=False):
    for dirname in i_file:
        del_dirs = os.path.join(i_path,dirname)
        if not os.listdir(del_dirs):
            os.rmdir(del_dirs)
            del_d += 1
a = f"\n已删除{del_no}个空文件"
b = f"已删除{del_copy}"
messagebox.showinfo("完成！",
                    f"空文件删除：{del_no}个\n"
                    f"重复文件删除：{del_copy}个\n"
                    f"空文件夹删除{del_d}个\n"
                    f"总计清除{del_no+del_copy+del_d}项\n")
